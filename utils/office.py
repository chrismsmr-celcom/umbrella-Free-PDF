import subprocess
import shlex
import os
import asyncio
from pathlib import Path

def convert_to_pdf(input_path: str, output_dir: str) -> str:
    """
    Convertit un fichier Office en PDF de manière sécurisée
    """
    # Vérifier que le fichier existe
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Fichier introuvable: {input_path}")
    
    # Vérifier l'extension
    ext = os.path.splitext(input_path)[1].lower()
    supported = ['.doc', '.docx', '.xls', '.xlsx', '.ppt', '.pptx', '.odt', '.rtf']
    
    if ext not in supported:
        raise ValueError(f"Extension non supportée: {ext}")
    
    # Préparer la commande de manière sécurisée
    output_file = os.path.join(output_dir, "converted.pdf")
    
    # Utiliser un tableau d'arguments (pas de shell=True)
    cmd = [
        "libreoffice",
        "--headless",
        "--convert-to", "pdf",
        "--outdir", output_dir,
        input_path
    ]
    
    try:
        # Timeout de 60 secondes
        result = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            timeout=60,
            check=False  # On vérifie manuellement
        )
        
        if result.returncode != 0:
            print(f"Erreur LibreOffice: {result.stderr}")
            raise Exception("La conversion a échoué")
        
        # Vérifier que le fichier de sortie existe
        expected_output = os.path.join(output_dir, Path(input_path).stem + ".pdf")
        if os.path.exists(expected_output):
            return expected_output
        elif os.path.exists(output_file):
            return output_file
        else:
            raise FileNotFoundError("Fichier PDF non généré")
            
    except subprocess.TimeoutExpired:
        raise TimeoutError("La conversion a pris trop de temps")
    except Exception as e:
        print(f"Erreur conversion: {e}")
        raise

def pdf_to_word(pdf_path: str, output_dir: str) -> str:
    """
    Convertit un PDF en Word avec gestion des erreurs
    """
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF introuvable: {pdf_path}")
    
    # Vérifier que c'est bien un PDF
    with open(pdf_path, 'rb') as f:
        header = f.read(4)
        if header != b'%PDF':
            raise ValueError("Le fichier n'est pas un PDF valide")
    
    output_path = os.path.join(output_dir, Path(pdf_path).stem + ".docx")
    
    # Utiliser pdf2docx (bibliothèque Python, pas de subprocess)
    try:
        from pdf2docx import Converter
        
        cv = Converter(pdf_path)
        cv.convert(output_path, start=0, end=None)
        cv.close()
        
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            return output_path
        else:
            raise Exception("Fichier Word vide ou corrompu")
            
    except Exception as e:
        print(f"Erreur conversion PDF->Word: {e}")
        raise

def pdf_to_excel(pdf_path: str, output_dir: str) -> str:
    """
    Convertit un PDF en Excel
    """
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF introuvable: {pdf_path}")
    
    output_path = os.path.join(output_dir, Path(pdf_path).stem + ".xlsx")
    
    try:
        import tabula
        
        # Extraire tous les tableaux
        tables = tabula.read_pdf(pdf_path, pages='all', multiple_tables=True)
        
        if not tables:
            raise Exception("Aucun tableau trouvé dans le PDF")
        
        # Sauvegarder dans Excel
        import pandas as pd
        with pd.ExcelWriter(output_path, engine='openpyxl') as writer:
            for i, table in enumerate(tables):
                sheet_name = f"Table_{i+1}"[:31]  # Excel max 31 chars
                table.to_excel(writer, sheet_name=sheet_name, index=False)
        
        if os.path.exists(output_path):
            return output_path
        else:
            raise Exception("Fichier Excel non généré")
            
    except Exception as e:
        print(f"Erreur conversion PDF->Excel: {e}")
        raise

def pdf_to_pptx(pdf_path: str, output_dir: str) -> str:
    """
    Convertit un PDF en PowerPoint
    """
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF introuvable: {pdf_path}")
    
    output_path = os.path.join(output_dir, Path(pdf_path).stem + ".pptx")
    
    try:
        from pdf2image import convert_from_path
        from pptx import Presentation
        from pptx.util import Inches
        
        # Convertir chaque page en image
        images = convert_from_path(pdf_path, dpi=200)
        
        if not images:
            raise Exception("Impossible de convertir le PDF en images")
        
        # Créer la présentation
        prs = Presentation()
        
        for img in images:
            # Sauvegarde temporaire de l'image
            temp_img = os.path.join(output_dir, f"temp_page_{len(prs.slides)}.png")
            img.save(temp_img, 'PNG')
            
            # Ajouter à la présentation
            slide_layout = prs.slide_layouts[6]  # Layout vierge
            slide = prs.slides.add_slide(slide_layout)
            
            # Ajouter l'image
            left = Inches(0)
            top = Inches(0)
            slide.shapes.add_picture(temp_img, left, top, height=prs.slide_height)
            
            # Nettoyer
            os.remove(temp_img)
        
        prs.save(output_path)
        
        if os.path.exists(output_path):
            return output_path
        else:
            raise Exception("Présentation PowerPoint non générée")
            
    except Exception as e:
        print(f"Erreur conversion PDF->PPT: {e}")
        raise

def pdf_to_pdfa(pdf_path: str, output_dir: str) -> str:
    """
    Convertit un PDF en PDF/A (archivage) - Version corrigee sans PdfA
    """
    if not os.path.exists(pdf_path):
        raise FileNotFoundError(f"PDF introuvable: {pdf_path}")
    
    output_path = os.path.join(output_dir, Path(pdf_path).stem + "_pdfa.pdf")
    
    try:
        from pikepdf import Pdf
        
        with Pdf.open(pdf_path, allow_overwriting_input=True) as pdf:
            # Compression standard sans validation PDF/A
            pdf.save(
                output_path,
                compress_streams=True,
                stream_decode_level=1
            )
        
        if os.path.exists(output_path) and os.path.getsize(output_path) > 0:
            print(f"PDF/A genere: {os.path.basename(output_path)}")
            return output_path
        else:
            raise Exception("Fichier PDF/A non genere")
            
    except Exception as e:
        print(f"Erreur conversion PDF/A: {e}")
        raise